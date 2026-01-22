from pptx import Presentation
import re
from loguru import logger
from src.utils.logo_fetcher import fetch_prospect_logo, fetch_intenthq_logo
from pptx.util import Inches

from pptx import Presentation
from pptx.util import Inches


def save_response_to_ppt(text: str, filename: str, prospect_url: str) -> str:
        """
        Save a presentation deck into a PowerPoint .pptx file.

        Parameters
        ----------
        deck : list[dict]
            Each dict represents a slide:
            {
                "title": "Slide title",
                "bullets": ["point 1", "point 2", ...]
            }
        filename : str
            Output PowerPoint file name.
        """
        logger.info("Saving response to PPT")
        prs = Presentation()

        # Fetch logos
        prospect_logo = fetch_prospect_logo(prospect_url)
        intenthq_logo = fetch_intenthq_logo()

        # Add logos to master slide
        # master = prs.slide_master

        logger.debug(f"Parsing presentation text for PPT {type(text)}")
        deck = parse_presentation(text)
        # Use a title + content layout for all slides
        title_and_content_layout = prs.slide_layouts[1]

        for slide_data in deck:
            slide = prs.slides.add_slide(title_and_content_layout)
            try:
                slide.shapes.add_picture(intenthq_logo, Inches(6), Inches(0.2), height=Inches(0.6))
            except Exception as e:
                logger.warning(f"Error adding IntentHQ logo to PPT: {e}")

            try:
                if prospect_logo:
                    slide.shapes.add_picture(prospect_logo, Inches(0.2), Inches(0.2), height=Inches(0.6))
            except Exception as e:
                logger.warning(f"Error adding Prospect logo to PPT: {e}")
            title_shape = slide.shapes.title
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame

            # Set title
            title_shape.text = slide_data.get("title", "")

            bullets = slide_data.get("bullets", [])

            if not bullets:
                text_frame.text = ""
                continue

            # First bullet
            text_frame.text = bullets[0]

            # Remaining bullets
            for bullet in bullets[1:]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0  # top-level bullet
        
        prs.save(filename)
        logger.info(f"PPT saved to {filename}")
        return filename

def parse_presentation(text: str):
    """
    Parse plain-text presentation into a deck structure:
    [
        {"title": "Slide title", "bullets": ["...", "..."], "description": "..."},
        ...
    ]
    """
    deck = []

    # Split on lines like: "Slide 1: Introduction"
    slide_pattern = re.compile(r"Slide\s+\d+:\s*(.+)")
    parts = slide_pattern.split(text)

    # parts[0] = preamble; then [slide_title, slide_body, slide_title, slide_body, ...]
    it = iter(parts[1:])
    for slide_heading, slide_body in zip(it, it):
        slide_dict = {
            "title": slide_heading.strip(),
            "bullets": [],
            "description": ""
        }

        lines = [l.strip() for l in slide_body.splitlines() if l.strip()]

        for line in lines:
            # Title line inside slide
            if line.lower().startswith("title:"):
                # Override slide title with this more specific title
                slide_dict["title"] = line.split(":", 1)[1].strip()
            # Bullets
            elif line.startswith("- "):
                bullet_text = line[2:].strip()
                slide_dict["bullets"].append(bullet_text)
            # Explanation line (single paragraph)
            elif line.lower().startswith("explanation:"):
                slide_dict["description"] = line.split(":", 1)[1].strip()
            else:
                # If there's extra text after Explanation, append it
                if slide_dict["description"]:
                    slide_dict["description"] += " " + line

        deck.append(slide_dict)

    return deck